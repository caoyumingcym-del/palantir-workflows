#!/usr/bin/env python3
"""
Build a deck for every experiment named by a set of sample manifests.

    python build_slides_batch.py                    # every sample_manifest*.csv here
    python build_slides_batch.py --dry-run          # say what it would do, do nothing
    python build_slides_batch.py --out-dir decks    # collect the decks in one folder
    python build_slides_batch.py --manifests a.csv b.csv --tables none

Where it looks
--------------
A manifest's ``output_path`` column is the experiment's output root, and the
pipeline writes its results to ``<output_path>/analysis_outputs``. So for each
manifest this reads::

    <output_path>/analysis_outputs/artifacts.json     <- what the deck is built from
    <output_path>/analysis_outputs/qc_report.html     <- run provenance and title

``output_path`` is a global column: every row of a manifest carries the same
value (the rows differ by DRAGEN run), so the file is read once and the distinct
values taken. A manifest naming several output paths is reported rather than
guessed at.

Anything else on the command line is passed straight through to
``build_slides.py``, so ``--tables none``, ``--pack width``, ``--appendix`` and
the rest work here exactly as they do for a single deck.

Why not a shell loop
--------------------
``for f in *.csv; do ...; done`` works until something is missing, and then it
either stops halfway or hides which experiments failed. This validates every
manifest before building anything, keeps going when one experiment has no
report yet, and finishes with a table saying what was built and what was not.
An explore-only run (``qc_explore.html``, no ``qc_report.html``) is called out
as such: those decks describe a QC pass, not a finished analysis.

Needs ``build_slides.py`` in the same folder.
"""
from __future__ import annotations

import argparse
import csv
import sys
import traceback
from dataclasses import dataclass, field
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

try:
    import build_slides
except ImportError:  # pragma: no cover
    sys.exit("build_slides.py must sit in the same folder as this script")

ANALYSIS_SUBDIR = "analysis_outputs"
REPORT_NAMES = ("qc_report.html", "qc_explore.html")


@dataclass
class Experiment:
    """One manifest, and what it points at."""

    manifest: Path
    output_path: Path | None = None
    artifacts: Path | None = None
    report: Path | None = None
    name: str = ""
    problems: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    # filled in after the run
    deck: Path | None = None
    slides: int = 0
    error: str = ""

    @property
    def ready(self) -> bool:
        return not self.problems


def read_output_paths(path: Path) -> tuple[list[str], list[str]]:
    """Distinct output_path values in a manifest, plus any problems reading it."""
    try:
        with path.open(newline="", encoding="utf-8-sig") as fh:
            rows = list(csv.DictReader(fh))
    except OSError as exc:
        return [], [f"cannot read: {exc}"]
    if not rows:
        return [], ["no data rows"]
    if "output_path" not in (rows[0].keys()):
        return [], ["no output_path column"]
    seen: list[str] = []
    for r in rows:
        value = (r.get("output_path") or "").strip()
        if value and value not in seen:
            seen.append(value)
    return seen, ([] if seen else ["output_path is empty in every row"])


def plan(manifests: list[Path], analysis_subdir: str) -> list[Experiment]:
    out: list[Experiment] = []
    for m in manifests:
        exp = Experiment(manifest=m)
        values, problems = read_output_paths(m)
        exp.problems.extend(problems)
        if values:
            if len(values) > 1:
                exp.problems.append(
                    "names several output_path values "
                    f"({', '.join(values[:3])}...); one manifest, one experiment"
                )
            else:
                exp.output_path = Path(values[0])
                exp.name = exp.output_path.name or m.stem
                analysis = exp.output_path / analysis_subdir
                if not exp.output_path.exists():
                    exp.problems.append(
                        f"output_path does not exist: {exp.output_path} "
                        "(wrong mount, or not this workspace?)"
                    )
                elif not analysis.is_dir():
                    # Tolerate a manifest whose output_path already *is* the
                    # analysis directory, or a layout with the results one level
                    # deeper than expected.
                    found = sorted(exp.output_path.rglob("artifacts.json"))
                    if found:
                        analysis = found[0].parent
                        exp.warnings.append(f"using {analysis}")
                    else:
                        exp.problems.append(
                            f"no {analysis_subdir}/ under {exp.output_path}"
                        )
                if not exp.problems:
                    art = analysis / "artifacts.json"
                    if art.exists():
                        exp.artifacts = art
                    else:
                        exp.problems.append(
                            f"no artifacts.json in {analysis} "
                            "(has the pipeline run for this experiment?)"
                        )
                    for name in REPORT_NAMES:
                        candidate = analysis / name
                        if candidate.exists():
                            exp.report = candidate
                            if name == "qc_explore.html":
                                exp.warnings.append(
                                    "explore run only: this deck describes a QC "
                                    "pass, not a finished analysis"
                                )
                            break
                    else:
                        exp.warnings.append(
                            "no qc_report.html: run provenance will be omitted"
                        )
        out.append(exp)
    return out


def deck_path(exp: Experiment, out_dir: Path | None) -> Path:
    if out_dir:
        return out_dir / f"{exp.name or exp.manifest.stem}_qc_slides.pptx"
    assert exp.artifacts is not None
    return exp.artifacts.parent / "qc_slides.pptx"


def run_one(exp: Experiment, out_dir: Path | None,
            passthrough: list[str]) -> None:
    """Build one deck in-process, surviving anything that goes wrong in it."""
    assert exp.artifacts is not None
    target = deck_path(exp, out_dir)
    if out_dir:
        out_dir.mkdir(parents=True, exist_ok=True)
    argv = [str(exp.artifacts), "--out", str(target), "--quiet"]
    if exp.report:
        argv += ["--from-report", str(exp.report)]
    argv += passthrough
    try:
        build_slides.main(argv)
    except SystemExit as exc:            # argparse or a guard inside the builder
        if exc.code not in (0, None):
            exp.error = f"build_slides exited with {exc.code}"
            return
    except Exception as exc:             # noqa: BLE001 - one bad run, not the batch
        exp.error = f"{type(exc).__name__}: {exc}"
        if VERBOSE_ERRORS:
            traceback.print_exc()
        return
    if not target.exists():
        exp.error = "no deck was written"
        return
    exp.deck = target
    try:
        from pptx import Presentation
        exp.slides = len(Presentation(str(target)).slides._sldIdLst)
    except Exception:                    # noqa: BLE001 - a count is not worth failing over
        exp.slides = 0


VERBOSE_ERRORS = False


def main(argv: list[str] | None = None) -> int:
    global VERBOSE_ERRORS
    p = argparse.ArgumentParser(
        description=__doc__.split("\n\n")[0],
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Unrecognised options are forwarded to build_slides.py, e.g. "
               "--tables none --pack width --appendix",
    )
    p.add_argument("--manifests", type=Path, nargs="*", metavar="CSV",
                   help="manifest files (default: sample_manifest*.csv beside "
                        "this script)")
    p.add_argument("--out-dir", type=Path,
                   help="write every deck here as <experiment>_qc_slides.pptx "
                        "(default: qc_slides.pptx beside each artifacts.json)")
    p.add_argument("--analysis-subdir", default=ANALYSIS_SUBDIR,
                   help=f"results directory under output_path "
                        f"(default: {ANALYSIS_SUBDIR})")
    p.add_argument("--dry-run", action="store_true",
                   help="show the plan and build nothing")
    p.add_argument("--stop-on-error", action="store_true",
                   help="stop at the first experiment that fails")
    p.add_argument("--traceback", action="store_true",
                   help="print a full traceback when an experiment fails")
    args, passthrough = p.parse_known_args(argv)
    VERBOSE_ERRORS = args.traceback

    manifests = args.manifests or sorted(HERE.glob("sample_manifest*.csv"))
    manifests = [Path(m) for m in manifests]
    if not manifests:
        print(f"No manifests found. Put sample_manifest*.csv next to "
              f"{HERE / 'build_slides_batch.py'}, or pass --manifests.",
              file=sys.stderr)
        return 2

    experiments = plan(manifests, args.analysis_subdir)

    # Two manifests writing one deck path would silently overwrite each other.
    if args.out_dir is None:
        by_target: dict[Path, list[Experiment]] = {}
        for e in experiments:
            if e.ready and e.artifacts:
                by_target.setdefault(deck_path(e, None), []).append(e)
        for target, group in by_target.items():
            if len(group) > 1:
                for e in group:
                    e.problems.append(
                        f"several manifests write to {target}; use --out-dir"
                    )

    width = max(len(e.manifest.name) for e in experiments)
    print(f"perturbseq slides — {len(experiments)} manifest(s)\n")
    for e in experiments:
        state = "ready" if e.ready else "SKIP "
        print(f"  {e.manifest.name:<{width}}  {state}  "
              f"{e.name or '?'}"
              + (f"  -> {deck_path(e, args.out_dir)}" if e.ready else ""))
        for w in e.warnings:
            print(f"  {'':<{width}}         note: {w}")
        for pr in e.problems:
            print(f"  {'':<{width}}         {pr}")
    print()

    runnable = [e for e in experiments if e.ready]
    if args.dry_run:
        print(f"--dry-run: {len(runnable)} deck(s) would be built.")
        return 0 if runnable else 1
    if not runnable:
        print("Nothing to build.", file=sys.stderr)
        return 1

    for e in runnable:
        print(f"building {e.name} ...", flush=True)
        run_one(e, args.out_dir, passthrough)
        if e.error:
            print(f"  failed: {e.error}", file=sys.stderr)
            if args.stop_on_error:
                break
        else:
            print(f"  {e.slides} slides -> {e.deck}")

    built = [e for e in runnable if e.deck and not e.error]
    failed = [e for e in experiments if e.error or not e.ready]
    print(f"\n{len(built)} deck(s) built, {len(failed)} not built.")
    for e in built:
        print(f"  ok    {e.name:<28} {e.slides:>3} slides  {e.deck}")
    for e in failed:
        reason = e.error or "; ".join(e.problems)
        print(f"  ----  {e.name or e.manifest.name:<28} {reason}")
    return 0 if not failed else 1


if __name__ == "__main__":
    raise SystemExit(main())
