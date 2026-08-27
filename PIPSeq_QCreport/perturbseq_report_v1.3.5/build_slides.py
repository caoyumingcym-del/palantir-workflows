#!/usr/bin/env python3
"""
Build a Google-Slides-ready deck from a perturb-seq ``artifacts.json``.

Why a .pptx and not the Slides API
----------------------------------
The Google Slides REST API cannot upload a local image: ``createImage`` takes a
publicly-fetchable URL, so an API-native builder would have to push 40 PNGs to a
world-readable bucket first, and carry OAuth credentials. A ``.pptx`` sidesteps
all of it. Google Slides imports PowerPoint natively:

* drag the file into Drive and open it with Google Slides -- one deck, done; or
* in an existing deck, **File > Import slides**, pick this file, tick the slides
  you want, and check "Keep original theme".

Either way text arrives as real text boxes, tables as real tables and figures as
native images -- everything stays editable and re-styleable, which is the point
of copy-pasting report slides into a group deck.

Why the registry and not the HTML
---------------------------------
``artifacts.json`` is the same source of truth ``report.py`` renders from, so the
deck inherits the section spine, the ordering, the captions, the metrics and the
"registered but not produced, and here is why" records for free. Parsing the
finished HTML would mean recovering that structure from base64 blobs.

This script deliberately has no dependency on the ``perturbseq_report`` package
-- ``artifacts.json`` serialises its own section spine, so a single file plus
``python-pptx`` is enough. (If the package *is* importable, the appendix picks up
the final review checklist from ``text.py``.)

Layout
------
The deck follows the report's own shape rather than inventing one:

  title  ->  contents  ->  [ per section: opener with blurb + metric cards +
             method notes,  then figure slides,  then table slides ]  ->  appendix

Under every slide title sits a two-line standfirst: **Shows** (what is on the
slide) and, where the report's own caption states one, **Good** / **Watch for** /
**How to read it** (how to tell whether it is fine). Both lines are extracted
from the registry's captions -- nothing about what "good" looks like is invented.

Figures are grouped by family where the report implies one (a run of
``<x> by <condition>`` panels becomes a single labelled grid), otherwise by the
declared or measured width. A deck carries no data tables: a table whose numbers
are already plotted is dropped, and one that is a label column plus numbers is
drawn as a bar chart instead (``--tables full`` restores real tables). Method
notes are merged where they repeat, curated down to the ones that report
something about this run, and budgeted per section. Every slide carries the run
label and the section title, so a slide pasted on its own into someone else's
deck still says what it is, and the full untruncated caption goes into the
speaker notes.

Nothing here is specific to one experiment. Sections, artifacts, captions and
severities all come from the registry, column roles in a table are inferred from
its values, and note curation is structural -- so a run with modalities this
script has never seen (hashtags, a new cross-check, a condition axis that did not
exist before) is laid out by the same rules.

Usage
-----
    python build_slides.py analysis/artifacts.json
    python build_slides.py analysis/ --out qc_deck.pptx --title "MDL-1898 QC"
    python build_slides.py analysis/artifacts.json --per-slide 3 --no-tables
    python build_slides.py analysis/artifacts.json --sections summary,guides

Requires: python-pptx (``pip install python-pptx``). Pillow is optional and only
used to downscale oversized figures.
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import html as _html
import json
import math
import re
import struct
import sys
import tempfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Iterable, Sequence

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:  # pragma: no cover
    sys.exit("python-pptx is required:  pip install python-pptx")


# ===========================================================================
# Look and feel -- the report's palette, so the deck and the HTML match
# ===========================================================================
C_INK        = RGBColor.from_string("1B1F24")
C_INK_SOFT   = RGBColor.from_string("4A5259")
C_INK_FAINT  = RGBColor.from_string("7A838C")
C_LINE       = RGBColor.from_string("E3E7EB")
C_LINE_SOFT  = RGBColor.from_string("F0F3F5")
C_ACCENT     = RGBColor.from_string("2F5D8A")
C_ACCENT_SOFT= RGBColor.from_string("EEF4FA")
C_GOOD       = RGBColor.from_string("2F7D4F")
C_GOOD_SOFT  = RGBColor.from_string("ECF7F0")
C_WARN       = RGBColor.from_string("9A6414")
C_WARN_SOFT  = RGBColor.from_string("FDF4E5")
C_POOR       = RGBColor.from_string("A3342F")
C_POOR_SOFT  = RGBColor.from_string("FDEEED")
C_WHITE      = RGBColor.from_string("FFFFFF")

LEVELS = {
    "good": (C_GOOD, C_GOOD_SOFT),
    "warn": (C_WARN, C_WARN_SOFT),
    "poor": (C_POOR, C_POOR_SOFT),
    "info": (C_ACCENT, C_ACCENT_SOFT),
}

# Arial is present in Google Slides, PowerPoint and Keynote; a theme font would
# be silently substituted on import and break the fitting maths below.
FONT = "Arial"

# Nothing in the deck is allowed below this. A QC report can use 8pt for a table
# of 25 columns because the reader is a foot from the screen; a slide cannot. The
# floor is enforced by making it the minimum of every fit_text call, so when text
# will not fit the layout gives it more room or the text is trimmed -- the font is
# never shrunk to compensate.
MIN_PT = 12.0

# Fallback spine, used only if artifacts.json predates the serialised sections.
DEFAULT_SECTIONS = [
    ("summary",       "Summary", "Headline numbers and the QC verdict for this experiment."),
    ("seq_qc",        "Sequencing QC", "Read-level metrics from the upstream alignment pipeline."),
    ("cell_qc",       "Per-cell QC & filtering", "Which cells were kept, on what basis, and what was lost."),
    ("transcriptome", "Transcriptome, embedding & clusters", "Normalisation, feature selection, embedding, clustering and markers."),
    ("guides",        "Guide assignment & performance", "Whether guides can be called, and how well."),
    ("perturbation",  "Perturbation effects", "Whether the perturbations did anything, and what."),
    ("hashtags",      "Hashtag performance", "Demultiplexing quality and sample composition."),
    ("comparability", "Condition comparability", "Whether the conditions being compared are comparable in the first place."),
    ("crosschecks",   "Cross-checks", "Independent measurements compared against each other."),
    ("appendix",      "Appendix", "Configuration, provenance and the review checklist."),
]


# ---------------------------------------------------------------------------
# Tables whose numbers are already plotted.
#
# A QC *report* can afford to print a figure and its underlying table; a deck
# cannot, and a 16-column lookup table spread over nine slides is worse than
# useless in a meeting. Two mechanisms remove them:
#
#   1. Title similarity (automatic, works for any pipeline): a table is dropped
#      when a figure in the same section carries essentially the same title --
#      "Guide performance by acoh" the table vs "Guide performance by acoh" the
#      figure, or "Knockdown per target" vs "Target-gene knockdown".
#   2. The list below, for the cases where the titles do not resemble each other
#      but the content is the same plot. Each entry names the figure that carries
#      the data, so the reasoning is auditable rather than a magic exclusion.
#
# Both are on by default; --keep-all-tables disables them, --keep-tables
# protects individual titles, --drop-tables removes more.
# ---------------------------------------------------------------------------
REDUNDANT_TABLES: tuple[tuple[str, str], ...] = (
    (r"^Depth measured from this h5ad$",
     "plotted in 'Depth before vs after downsampling'"),
    (r"^Cells rejected per QC gate$",
     "plotted in 'Cell retention'"),
    (r"^Sequencing metrics\b",
     "plotted in 'Sequencing metrics per library'"),
    (r"^Knockdown per target$",
     "plotted in 'Target-gene knockdown'"),
    (r"^Guide ID to target mapping$",
     "a 16-column library lookup; belongs in the CSV, not a deck"),
)

# ---------------------------------------------------------------------------
# Which method notes earn a slide.
#
# The rules are structural, not a list of titles, because the notes a run emits
# depend on its modalities: a hashtag experiment produces demultiplexing notes a
# guide-only run never will, and the appendix grows a block per input file. So:
#
#   * notes repeated under one title are merged into a single bulleted block
#     (eight "Hashtag recovery from DRAGEN output" notes are one block, not
#     eight);
#   * a note that reports no numbers at all is explanatory prose rather than a
#     finding about this run -- the slide standfirst already carries the "how to
#     read it" job -- so it is dropped;
#   * what survives is ranked poor > warn > info, then by whether it quantifies
#     something, and each section keeps the top few.
#
# Nothing here knows what a hashtag or a guide is, so a run with modalities this
# script has never seen is curated the same way.
# ---------------------------------------------------------------------------
def _is_inventory(body: str) -> bool:
    """True for notes that are mostly a list of names.

    "39 metric name(s) ... were ignored: Feature reads with non-matching
    barcodes, Filtered ambiguously matching reads, ..." quotes a number, so the
    numeric test keeps it, but a slide of column names nobody will read is not a
    finding. Long comma-separated runs of short items are the signature.
    """
    text = plain(body)
    items = [x.strip() for x in text.split(",") if x.strip()]
    if len(items) < 8:
        return False
    return sum(len(x) for x in items) / len(items) < 60


EXPLANATORY_TITLE = re.compile(
    r"^(about|reading the|method|two definitions|what space|how to)\b", re.I)
# Keys the pipeline uses for content that must survive curation regardless.
ALWAYS_KEEP_NOTE_KEYS = ("checklist",)


def _has_numbers(body: str) -> bool:
    return bool(re.search(r"\d", plain(body)))


def curate_notes(notes: Sequence[Art], max_per_section: int,
                 dropped: list[tuple[str, str]]) -> list[Art]:
    """Merge repeated notes, drop pure prose, keep the most informative few."""
    order: list[str] = []
    grouped: dict[str, list[Art]] = {}
    for a in notes:
        t = plain(a.title)
        if t not in grouped:
            grouped[t] = []
            order.append(t)
        grouped[t].append(a)

    merged: list[Art] = []
    for t in order:
        group = grouped[t]
        levels = [str(a.data.get("level", "info")) for a in group]
        level = ("poor" if "poor" in levels
                 else "warn" if "warn" in levels else "info")
        bodies = [plain(a.data.get("body") or a.caption) for a in group]
        bodies = [b for b in bodies if b]
        body = (bodies[0] if len(bodies) == 1
                else "\n".join(f"•  {b}" for b in bodies))
        merged.append(Art(
            kind="note", section=group[0].section, key=group[0].key, title=t,
            caption="", order=group[0].order, width="full",
            data={"body": body, "level": level,
                  "merged": len(group) if len(group) > 1 else 0},
        ))

    keep: list[Art] = []
    for a in merged:
        body = a.data.get("body", "")
        level = str(a.data.get("level", "info"))
        if a.key in ALWAYS_KEEP_NOTE_KEYS or level in ("warn", "poor"):
            keep.append(a)
            continue
        if not _has_numbers(body) or EXPLANATORY_TITLE.match(plain(a.title)):
            dropped.append((plain(a.title),
                            "explanatory prose, not a finding about this run"))
            continue
        if _is_inventory(body):
            dropped.append((plain(a.title),
                            "a list of names rather than a finding"))
            continue
        keep.append(a)

    if max_per_section and len(keep) > max_per_section:
        ranked = sorted(
            keep,
            key=lambda a: note_rank(str(a.data.get("level", "info")),
                                    a.data.get("body", ""), a.order),
        )
        chosen = {id(a) for a in ranked[:max_per_section]}
        chosen |= {id(a) for a in keep if a.key in ALWAYS_KEEP_NOTE_KEYS}
        for a in keep:
            if id(a) not in chosen:
                dropped.append((plain(a.title),
                                f"over the {max_per_section}-note budget for "
                                f"its section"))
        keep = [a for a in keep if id(a) in chosen]
    return keep


def note_rank(level: str, body: str, order: int) -> tuple:
    """Sort key: severity first, then whether it quantifies, then report order."""
    severity = {"poor": 3, "warn": 2}.get(level, 0)
    return (-severity, not _has_numbers(body), order)


_WORD = re.compile(r"[a-z0-9]+")
_STOP = {"by", "per", "the", "of", "in", "a", "and", "from", "this", "for"}


def _tokens(title: str) -> set[str]:
    words = []
    for w in _WORD.findall(plain(title).lower()):
        if w in _STOP:
            continue
        # crude singularisation so "Call rates vs threshold" (figure) matches
        # "Threshold and call rate by normalisation" (table)
        if len(w) > 3 and w.endswith("s") and not w.endswith("ss"):
            w = w[:-1]
        words.append(w)
    return set(words)


def table_is_plotted(table_title: str, figure_titles: Sequence[str],
                     threshold: float = 0.5) -> str | None:
    """Return the figure that already shows this table, if any.

    Deliberately conservative. A short table title is *not* enough on its own:
    "QC thresholds applied" is a subset of "QC metrics and thresholds (before
    filtering)" by words, but the table's ``source`` column -- whether a cut-off
    was derived automatically or set by a human -- appears in no figure. So the
    subset shortcut requires at least three meaningful words, and everything
    else has to clear the Jaccard threshold. The best-scoring figure is
    reported, not the first one that happens to match.
    """
    t = _tokens(table_title)
    if not t:
        return None
    best, best_score = None, 0.0
    for f in figure_titles:
        g = _tokens(f)
        if not g:
            continue
        inter = len(t & g)
        if not inter:
            continue
        score = inter / len(t | g)
        if (len(t) >= 3 and t <= g) or len(g) >= 3 and g <= t:
            score = max(score, 1.0)
        if score >= threshold and score > best_score:
            best, best_score = plain(f), score
    return best


# ===========================================================================
# Registry loading
# ===========================================================================
@dataclass
class Art:
    """One artifact, as recorded by ``artifacts.py``."""

    kind: str
    section: str
    key: str
    title: str
    path: str | None = None
    caption: str = ""
    order: int = 100
    width: str = "half"
    data: dict[str, Any] = field(default_factory=dict)
    skipped_reason: str | None = None

    @property
    def produced(self) -> bool:
        return self.skipped_reason is None

    def resolve(self, root: Path) -> Path | None:
        if not self.path:
            return None
        p = Path(self.path)
        return p if p.is_absolute() else (root / p)


@dataclass
class Registry:
    root: Path
    sections: list[tuple[str, str, str]]
    artifacts: list[Art]

    @classmethod
    def load(cls, path: Path) -> "Registry":
        payload = json.loads(path.read_text(encoding="utf-8"))
        # Root comes from where the file actually is, not from the absolute path
        # recorded at write time, so a copied analysis directory still resolves.
        root = path.parent
        raw_sections = payload.get("sections") or []
        sections = (
            [(s["key"], s["title"], s.get("blurb", "")) for s in raw_sections]
            or list(DEFAULT_SECTIONS)
        )
        known = set(Art.__dataclass_fields__)
        arts = [
            Art(**{k: v for k, v in item.items() if k in known})
            for item in payload.get("artifacts", [])
        ]
        return cls(root=root, sections=sections, artifacts=arts)

    def by_section(self) -> dict[str, list[Art]]:
        order = {k: i for i, (k, _, _) in enumerate(self.sections)}
        out: dict[str, list[Art]] = {}
        for a in sorted(
            self.artifacts,
            key=lambda x: (order.get(x.section, 999), x.order, x.key),
        ):
            out.setdefault(a.section, []).append(a)
        return out


# ===========================================================================
# Text: the report's captions and notes are HTML; slides need plain prose
# ===========================================================================
_TAG = re.compile(r"<[^>]+>")
_WS = re.compile(r"[ \t\r\f\v]+")


_BLANKS = re.compile(r"\n{2,}")
_BREAK = re.compile(r"<\s*(br|/p|/div|/h[1-6]|/tr|/li)\s*/?\s*>", re.I)
_ITEM = re.compile(r"<\s*li[^>]*>", re.I)


def plain(value: object) -> str:
    """HTML caption/note body -> plain text suitable for a text box.

    The report's note bodies are authored HTML (``text.py`` emits ``<p>`` and
    ``<ul>``), so block boundaries become line breaks and list items become
    bullets instead of being flattened into one run-on paragraph.
    """
    if value is None:
        return ""
    s = _ITEM.sub("\n\u2022  ", str(value))
    s = _BREAK.sub("\n", s)
    s = _TAG.sub(" ", s)
    s = _html.unescape(s).replace("\xa0", " ")
    lines = [_WS.sub(" ", ln).strip() for ln in s.split("\n")]
    out = "\n".join(ln for ln in lines if ln)
    return _BLANKS.sub("\n", out).strip()


_SENT = re.compile(r"(?<=[.!?])\s+(?=[A-Z(“])")


def first_sentences(text: str, n: int) -> str:
    if n <= 0 or not text:
        return ""
    parts = _SENT.split(text)
    kept = " ".join(parts[:n]).strip()
    return kept


def sentences(text: str) -> list[str]:
    return [x.strip() for x in _SENT.split(plain(text).replace("\n", " ")) if x.strip()]


def fit_sentences(text: str, width_in: float, height_in: float,
                  size: float) -> tuple[str, bool]:
    """The longest run of *complete* sentences that fits. Never cuts mid-clause.

    A caption ending in an ellipsis is not a description, it is a tease. So the
    deck drops whole sentences rather than characters, and whatever it drops
    goes into the speaker notes.
    """
    parts = sentences(text)
    if not parts:
        return "", False
    max_lines = max(1, int(height_in * 72.0 / (size * _LINE)))
    kept: list[str] = []
    for part in parts:
        trial = " ".join(kept + [part])
        if kept and _line_count(trial, width_in, size) > max_lines:
            break
        kept.append(part)
    if not kept:
        kept = [parts[0]]
    return " ".join(kept), len(kept) < len(parts)


def truncate(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    cut = text[: max(0, limit - 1)]
    sp = cut.rfind(" ")
    if sp > limit * 0.6:
        cut = cut[:sp]
    return cut.rstrip(" ,;:.") + "…"


# ---------------------------------------------------------------------------
# The standfirst under each slide title: what the slide shows, and how to tell
# whether it is good.
#
# The report's captions already contain both -- text.py was written to explain
# how to read each panel -- so this extracts rather than invents. The first
# sentence is almost always descriptive ("UMAP of QC-passing cells coloured
# by..."), and somewhere later there is usually a sentence stating the criterion
# ("Clean single-guide cells sit in the top-right corner", "A cluster drawn
# almost entirely from one sample is the clearest signature of a batch effect").
# Nothing is fabricated: when no such sentence exists, the line is simply
# omitted rather than filled with a guess about what good looks like.
# ---------------------------------------------------------------------------
_GOOD_CUE = re.compile(
    r"\b(good|healthy|clean|expect(?:ed)?|should|ideal|want|well[- ]behaved|"
    r"reassuring|acceptable|robust|plateau|consistent|sit[s]? (?:at|in|near)|"
    r"centred|centered)\b", re.I)
_RISK_CUE = re.compile(
    r"\b(bad|worry\w*|problem\w*|red flag|signature of|suspicious|fail\w*|"
    r"concern\w*|broken|artefact|artifact|batch effect|contamina\w*|"
    r"underpowered|cannot|misleading|skew\w*)\b", re.I)
# Third tier: no explicit good/bad word, but the sentence still tells you how to
# read the panel ("... is diagnostic", "a steep curve means ...").
_READ_CUE = re.compile(
    r"\b(means|indicates|suggests|diagnostic|defensible|disproportion\w*|"
    r"tells you|look for|worth checking|the panel to read|sign\b|"
    r"cannot support|barely moves|sensitive to)\b", re.I)
# Sentences opening with a connective or bare pronoun refer back to something
# that will not be on the slide, so they read as non-sequiturs in isolation.
_DANGLING = re.compile(
    r"^(so|if|then|it|they|this|these|that|those|otherwise|hence|thus|and|but|"
    r"because|there|here|both|either|neither)\b", re.I)


def common_sentences(captions: Sequence[str]) -> str:
    """The sentences every caption in a family shares, in their original order.

    Members of a family differ only in the condition they name ("...for each
    level of fixation" vs "...of gRNA_method"), so the sentences they have in
    common are exactly the part that is true of the whole grid.
    """
    texts = [plain(c) for c in captions if plain(c)]
    if not texts:
        return ""
    if len(texts) == 1:
        return texts[0]
    per_caption = [
        [x.strip() for x in _SENT.split(t) if x.strip()] for t in texts
    ]
    shared = set(per_caption[0])
    for group in per_caption[1:]:
        shared &= set(group)
    return " ".join(x for x in per_caption[0] if x in shared)


def caption_brief(caption: str, max_chars: int = 0) -> tuple[str, str, str]:
    """(what the slide shows, the criterion, label for the criterion)."""
    text = plain(caption)
    if not text:
        return "", "", ""
    sents = [x.strip() for x in _SENT.split(text.replace("\n", " ")) if x.strip()]
    if not sents:
        return "", "", ""
    shows = sents[0]
    # "What this shows." / "How the DEGs are selected." are headings, not prose.
    if len(shows) < 42 and len(sents) > 1:
        shows = f"{shows} {sents[1]}".strip()
        rest = sents[2:]
    else:
        rest = sents[1:]

    best, best_score, has_good = "", 0.0, False
    for sent in rest:
        if _DANGLING.match(sent):
            continue
        good_hits = len(_GOOD_CUE.findall(sent))
        risk_hits = len(_RISK_CUE.findall(sent))
        read_hits = len(_READ_CUE.findall(sent))
        if not (good_hits or risk_hits or read_hits):
            continue
        score = (3.0 * good_hits + 2.0 * risk_hits + 1.0 * read_hits
                 - len(sent) / 400.0)
        if score > best_score:
            best, best_score, has_good = sent, score, bool(good_hits)
    label = ("Good" if has_good else "Watch for") if best else ""
    if not best:
        # No sentence advertises a criterion. These captions still tend to end
        # on the interpretive punchline ("...adding more reads will not fix
        # it."), so fall back to the closing sentence rather than inventing a
        # standard of "good" the report never claimed.
        for sent in reversed(rest):
            if len(sent) >= 40 and not _DANGLING.match(sent):
                best, label = sent, "How to read it"
                break
    if best and best[:60] == shows[:60]:
        best, label = "", ""
    return shows, best, label


# Average glyph advance for Arial as a fraction of point size. Empirical, and
# only used to pick a font size that will not overflow -- being slightly
# pessimistic here is much better than a text box that spills off the slide.
_CHAR_EM = 0.53
_LINE = 1.24


def _line_count(text: str, width_in: float, size_pt: float) -> int:
    per_line = max(6, int((width_in * 72.0) / (_CHAR_EM * size_pt)))
    return sum(max(1, math.ceil(len(ln) / per_line)) for ln in text.split("\n"))


def fit_text(
    text: str, width_in: float, height_in: float,
    max_pt: float, min_pt: float,
) -> tuple[str, float]:
    """Largest size in [min_pt, max_pt] that fits; truncates if even min won't."""
    if not text:
        return "", max_pt
    pt = max_pt
    while pt >= min_pt:
        if _line_count(text, width_in, pt) * pt * _LINE / 72.0 <= height_in:
            return text, pt
        pt -= 0.5
    per_line = max(6, int((width_in * 72.0) / (_CHAR_EM * min_pt)))
    max_lines = max(1, int(height_in * 72.0 / (min_pt * _LINE)))
    return truncate(text, per_line * max_lines), min_pt


# ===========================================================================
# Images
# ===========================================================================
def png_size(path: Path) -> tuple[int, int] | None:
    """Width/height from the PNG header, without pulling in Pillow."""
    try:
        with path.open("rb") as fh:
            head = fh.read(26)
        if head[:8] != b"\x89PNG\r\n\x1a\n" or head[12:16] != b"IHDR":
            return None
        w, h = struct.unpack(">II", head[16:24])
        return int(w), int(h)
    except OSError:
        return None


def image_size(path: Path) -> tuple[int, int]:
    got = png_size(path)
    if got:
        return got
    try:
        from PIL import Image  # optional
        with Image.open(path) as im:
            return im.size
    except Exception:
        return (1600, 1200)  # assume 4:3 rather than crash


def file_digest(path: Path) -> str | None:
    try:
        return hashlib.md5(path.read_bytes()).hexdigest()
    except OSError:
        return None


def maybe_downscale(path: Path, max_px: int, cache: Path) -> Path:
    """Shrink very large figures so the deck stays inside Drive's limits."""
    if max_px <= 0:
        return path
    w, h = image_size(path)
    if max(w, h) <= max_px:
        return path
    try:
        from PIL import Image
    except ImportError:
        return path
    try:
        cache.mkdir(parents=True, exist_ok=True)
        target = cache / f"{path.stem}_{max_px}{path.suffix or '.png'}"
        if not target.exists():
            with Image.open(path) as im:
                im.thumbnail((max_px, max_px), Image.LANCZOS)
                im.save(target)
        return target
    except Exception:
        return path


# ===========================================================================
# Value formatting -- mirrors report._format_cell so numbers agree with the HTML
# ===========================================================================
_NUMERIC = re.compile(r"^[-+]?(\d+\.?\d*|\.\d+)([eE][-+]?\d+)?$")


def fmt_value(value: Any) -> tuple[str, bool]:
    if value is None:
        return "–", False
    if isinstance(value, str) and _NUMERIC.match(value.strip()):
        # Values read back from a CSV arrive as strings; format and align them
        # as numbers so a CSV-backed table matches an inline one.
        txt = value.strip()
        value = int(txt) if _NUMERIC.match(txt) and txt.lstrip("+-").isdigit() else float(txt)
    if isinstance(value, bool):
        return ("yes" if value else "no"), False
    if isinstance(value, int):
        return f"{value:,}", True
    if isinstance(value, float):
        if value != value:
            return "–", True
        if value == int(value) and abs(value) < 1e15:
            return f"{int(value):,}", True
        if abs(value) < 1e-4 and value != 0:
            return f"{value:.2e}", True
        return f"{value:,.4g}", True
    return plain(value), False


# ===========================================================================
# Optional: recover what the HTML report knows and artifacts.json does not
# ===========================================================================
def parse_report_html(path: Path) -> tuple[str | None, list[dict[str, str]]]:
    """Lift the title and the appendix provenance table out of a built report.

    ``pipeline.py`` assembles the provenance dict at render time and never
    serialises it, so it is the one thing in the HTML that is not in
    ``artifacts.json``. Rather than re-derive it (and risk disagreeing with the
    report), read it back out of the report itself.
    """
    text = path.read_text(encoding="utf-8", errors="replace")
    m = re.search(r"<title>(.*?)</title>", text, re.S)
    title = plain(m.group(1)) if m else None
    rows: list[dict[str, str]] = []
    idx = text.rfind('id="sec-appendix"')
    if idx != -1:
        for tr in re.findall(r"<tr>(.*?)</tr>", text[idx:], re.S):
            cells = re.findall(r"<td[^>]*>(.*?)</td>", tr, re.S)
            if len(cells) == 2:
                rows.append({"item": plain(cells[0]), "value": plain(cells[1])})
    return title, rows


# ===========================================================================
# Tables as charts
#
# A deck should not carry a 50-row table, but dropping one silently loses what
# it measured. Where a table is a label column plus numbers, it is drawn as a
# bar chart instead: the shape of the distribution is the part that belongs on a
# slide, and the exact numbers are in the CSV the pipeline already wrote.
#
# Column roles are inferred, not configured, so this works on a table this
# script has never seen: the first mostly-non-numeric column labels the bars,
# and the value is the numeric column whose name reads like a quantity. A table
# with several unnamed numeric columns (a cross-tab) is drawn stacked.
# ===========================================================================
# Tried in order: a proportion beats a count, a count beats a raw total, and a
# derived statistic is the last resort. Ordering matters more than membership --
# "Cells per target" has n_guides, total_umis and n_cells_assigned, and only one
# of those is what the table is about.
VALUE_CUES = (
    re.compile(r"(pct|percent|frac|proportion|rate)", re.I),
    re.compile(r"(cells|count|n_obs)", re.I),
    re.compile(r"(umis|reads|total|depth|size)", re.I),
    re.compile(r"(knockdown|shift|score|separation|median|mean|value)", re.I),
)


def _as_number(value: Any) -> float | None:
    if isinstance(value, bool):
        return None
    if isinstance(value, (int, float)):
        return None if value != value else float(value)
    if isinstance(value, str):
        txt = value.strip().replace(",", "")
        if _NUMERIC.match(txt):
            return float(txt)
    return None


def table_columns_by_role(cols: Sequence[str], rows: Sequence[dict]
                          ) -> tuple[list[str], list[str]]:
    """(label columns, numeric columns) inferred from the values themselves.

    An all-integer first column with no repeats is an identifier, not a
    measurement -- a cluster table's ``cluster`` column would otherwise be
    charted as if 0..23 were data, and the long ``description`` column would end
    up labelling the bars.
    """
    numeric: list[str] = []
    labels: list[tuple[int, str]] = []
    for i, c in enumerate(cols):
        vals = [_as_number(r.get(c)) for r in rows]
        hits = sum(v is not None for v in vals)
        raw = [plain(r.get(c)) for r in rows]
        distinct = len(set(raw))
        if rows and hits >= 0.8 * len(rows):
            is_id = (i == 0 and distinct == len(rows)
                     and all(v is not None and v == int(v) for v in vals))
            (labels.append((distinct, c)) if is_id else numeric.append(c))
        elif rows and hits <= 0.2 * len(rows):
            labels.append((distinct, c))
    # Most distinct labels first: that is the column the rows are really keyed by.
    ordered = [c for _, c in sorted(labels, key=lambda x: -x[0])]
    return ordered, numeric


# Charts are authored at the size they will be placed (11.5in wide inside a
# 12.23in content column), so 12pt in the PNG is 12pt on the slide. Rendering a
# 9.5in figure and letting fit-contain shrink it to 5in would put the tick
# labels at 6pt -- below the floor everything else in the deck respects.
CHART_W_IN = 11.5
CHART_MAX_H_IN = 4.3


def chart_table(a: Art, cols: Sequence[str], rows: Sequence[dict],
                out_dir: Path, max_bars: int = 12) -> tuple[Path, str] | None:
    """Render a table as a bar chart. Returns (png, what was charted)."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception:
        return None
    if len(rows) < 3:
        return None
    label_cols, numeric = table_columns_by_role(cols, rows)
    if not label_cols or not numeric:
        return None
    label_col = label_cols[0]

    series: list[str] = []
    for cue in VALUE_CUES:
        hit = [c for c in numeric if cue.search(c)]
        if hit:
            series = [hit[0]]
            break
    if not series:
        series = numeric if len(numeric) > 1 else numeric[:1]
    if len(series) > 6:
        series = series[:6]

    def bar_label(r: dict) -> str:
        text = plain(r.get(label_col))
        # Rows keyed by a pair of columns ("CLR" x "hash.A") need both, or the
        # chart shows the same label several times.
        if len(label_cols) > 1 and len(set(
                plain(x.get(label_col)) for x in rows)) < len(rows):
            second = plain(r.get(label_cols[1]))
            if second:
                text = f"{text} · {second}"
        return text

    data = []
    for r in rows:
        vals = [(_as_number(r.get(c)) or 0.0) for c in series]
        if any(v for v in vals):
            data.append((bar_label(r), vals))
    if len(data) < 3:
        return None
    data.sort(key=lambda x: sum(x[1]), reverse=True)
    clipped = len(data) > max_bars
    data = data[:max_bars]

    if len(series) == 1:
        positive = [v[0] for _, v in data if v[0] > 0]
        # Mixed units in one column (a table of assorted thresholds) make a bar
        # chart that says nothing; leave those alone.
        if positive and max(positive) / min(positive) > 1e3:
            return None

    labels = [truncate(lbl, 34) or "—" for lbl, _ in data]
    height = min(CHART_MAX_H_IN, max(2.4, 0.26 * len(data) + 1.1))
    fig, ax = plt.subplots(figsize=(CHART_W_IN, height), dpi=200)
    y = list(range(len(data)))[::-1]
    palette = ["#2f5d8a", "#7aa5c9", "#a3342f", "#2f7d4f", "#9a6414", "#7a838c"]
    left = [0.0] * len(data)
    for si, col in enumerate(series):
        vals = [v[si] for _, v in data]
        ax.barh(y, vals, left=left, color=palette[si % len(palette)],
                label=plain(col), height=0.72)
        left = [l + v for l, v in zip(left, vals)]
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=MIN_PT)
    ax.set_xlabel(plain(series[0]) if len(series) == 1 else "value",
                  fontsize=MIN_PT)
    ax.tick_params(axis="x", labelsize=MIN_PT)
    # No chart title: the slide already labels the figure with the same words.
    if len(series) > 1:
        ax.legend(fontsize=MIN_PT, frameon=False, ncol=min(len(series), 3))
    what = ", ".join(plain(c) for c in series) + f" per {plain(label_col)}"
    if clipped:
        what = f"Top {len(data)} of {len(rows):,} rows by {what}"
        # Stamped on the image as well, so a slide pasted somewhere else still
        # says it is a partial view.
        ax.set_title(f"top {len(data)} of {len(rows):,} rows",
                     fontsize=MIN_PT, loc="right", color="#7a838c")
    else:
        what = f"All {len(rows):,} rows: {what}"
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    ax.grid(axis="x", color="#e3e7eb", linewidth=0.8)
    ax.set_axisbelow(True)
    fig.tight_layout()
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / f"chart_{a.section}_{a.key}.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)

    return path, what


def tables_as_figures(tables: Sequence[Art], root: Path, out_dir: Path,
                      max_rows: int, charted: list[tuple[str, str]],
                      skipped: list[tuple[str, str]]) -> list[Art]:
    """Turn each chartable table into a figure artifact the deck can lay out."""
    made: list[Art] = []
    for a in tables:
        cols, rows, _total = table_rows(a, root, max_rows * 4)
        if not cols or not rows:
            skipped.append((plain(a.title), "no rows to chart"))
            continue
        result = chart_table(a, cols, rows, out_dir)
        if result is None:
            skipped.append((plain(a.title), "not chartable (no label + value "
                                            "columns, or mixed units)"))
            continue
        path, what = result
        cap = plain(a.caption)
        # Scope first: the standfirst quotes the opening sentence, and a reader
        # has to know this is a chart of part of a table before anything else.
        made.append(Art(
            kind="figure", section=a.section, key=f"chart_{a.key}",
            title=plain(a.title),
            path=str(path), order=a.order + 500, width="full",
            caption=(f"{what}, charted from the table of the same name; the "
                     f"full table is in the CSV alongside the report. {cap}"),
            data={"charted": True},
        ))
        charted.append((plain(a.title), what))
    return made


# ===========================================================================
# Slide geometry
# ===========================================================================
@dataclass
class Geometry:
    width: float
    height: float
    margin: float = 0.55
    header_h: float = 0.95
    footer_h: float = 0.42

    @property
    def content_left(self) -> float:
        return self.margin

    @property
    def content_width(self) -> float:
        return self.width - 2 * self.margin

    @property
    def content_top(self) -> float:
        return self.margin + self.header_h

    @property
    def content_height(self) -> float:
        return self.height - self.content_top - self.footer_h - 0.18


class Deck:
    """Thin wrapper over python-pptx that keeps every slide on one grid."""

    def __init__(self, aspect: str, label: str, generated: str):
        self.prs = Presentation()
        if aspect == "4:3":
            w, h = 10.0, 7.5
        else:
            w, h = 13.333, 7.5
        self.prs.slide_width = Inches(w)
        self.prs.slide_height = Inches(h)
        self.geo = Geometry(w, h)
        self.label = label
        self.generated = generated
        self.n = 0

    # -------------------------------------------------------------- plumbing
    def blank(self):
        # Layout 6 of the default template is the blank one: nothing inherited,
        # so "Keep original theme" on import does not drag in stray placeholders.
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def text(
        self, slide, txt: str, left: float, top: float, width: float, height: float,
        size: float = 12, bold: bool = False, color: RGBColor = C_INK,
        align=PP_ALIGN.LEFT, italic: bool = False, caps: bool = False,
        spacing: float | None = None, anchor=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        first = True
        for line in (txt or "").split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.line_spacing = _LINE
            if spacing is not None:
                p.space_after = Pt(spacing)
            run = p.add_run()
            run.text = line.upper() if caps else line
            f = run.font
            f.name = FONT
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.color.rgb = color
        return box

    def rich(self, slide, paragraphs, left: float, top: float, width: float,
             height: float, size: float):
        """Text box whose paragraphs mix formatting: [(text, bold, colour), ...]."""
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        for i, segments in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = _LINE
            for text, bold, colour in segments:
                run = p.add_run()
                run.text = text
                run.font.name = FONT
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = colour
        return box

    def brief(self, slide, shows: str, verdict: str = "", label: str = "",
              size: float = MIN_PT, top: float | None = None) -> float:
        """Standfirst under the title. Returns the height it used."""
        g = self.geo
        paras, plain_lines = [], []
        if shows:
            paras.append([("Shows  ", True, C_ACCENT), (shows, False, C_INK_SOFT)])
            plain_lines.append(f"Shows  {shows}")
        if verdict:
            colour = C_GOOD if label == "Good" else C_WARN
            paras.append([(f"{label}  ", True, colour), (verdict, False, C_INK_SOFT)])
            plain_lines.append(f"{label}  {verdict}")
        if not paras:
            return 0.0
        lines = sum(_line_count(t, g.content_width, size) for t in plain_lines)
        h = lines * size * _LINE / 72.0 + 0.04
        self.rich(slide, paras, g.content_left,
                  (g.content_top if top is None else top) - 0.04,
                  g.content_width, h, size)
        return h + 0.12

    def rule(self, slide, left: float, top: float, width: float,
             color: RGBColor = C_LINE, thickness: float = 0.012):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
            Inches(width), Inches(thickness),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def panel(self, slide, left: float, top: float, width: float, height: float,
              fill: RGBColor, line: RGBColor | None = None,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE):
        shp = slide.shapes.add_shape(
            shape, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(0.75)
        shp.shadow.inherit = False
        try:  # keep the corner radius subtle rather than pill-shaped
            shp.adjustments[0] = 0.06
        except (IndexError, ValueError):
            pass
        return shp

    def picture(self, slide, path: Path, left: float, top: float,
                width: float, height: float) -> float:
        """Fit-contain a figure in the box; returns the height actually used."""
        iw, ih = image_size(path)
        scale = min(width / iw, height / ih)
        w, h = iw * scale, ih * scale
        pic = slide.shapes.add_picture(
            str(path), Inches(left + (width - w) / 2), Inches(top),
            Inches(w), Inches(h),
        )
        pic.line.color.rgb = C_LINE
        pic.line.width = Pt(0.5)
        return h

    # ----------------------------------------------------------- slide chrome
    def new(self, title: str, eyebrow: str | None = None, rule: bool = True):
        slide = self.blank()
        g = self.geo
        self.n += 1
        self.text(
            slide, eyebrow or self.label, g.margin, g.margin - 0.06,
            g.content_width, 0.24, size=MIN_PT, bold=True, color=C_ACCENT,
            caps=True,
        )
        txt, size = fit_text(title, g.content_width, 0.46, 23, 16)
        self.text(slide, txt, g.margin, g.margin + 0.24, g.content_width, 0.5,
                  size=size, bold=True, color=C_INK)
        if rule:
            self.rule(slide, g.margin, g.margin + 0.78, g.content_width)
        self.footer(slide)
        return slide

    def footer(self, slide):
        g = self.geo
        y = self.geo.height - 0.42
        self.rule(slide, g.margin, y - 0.1, g.content_width, C_LINE_SOFT)
        self.text(slide, f"Perturb-seq QC · {self.generated}", g.margin, y,
                  g.content_width - 0.7, 0.24, size=MIN_PT, color=C_INK_FAINT)
        self.text(slide, str(self.n), g.width - g.margin - 0.6, y, 0.6, 0.24,
                  size=MIN_PT, color=C_INK_FAINT, align=PP_ALIGN.RIGHT)

    def notes(self, slide, text: str):
        if text:
            slide.notes_slide.notes_text_frame.text = text

    def save(self, path: Path) -> Path:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return path


# ===========================================================================
# Slide builders
# ===========================================================================
def title_slide(deck: Deck, title: str, subtitle: str, meta: Sequence[str]):
    slide = deck.blank()
    deck.n += 1
    g = deck.geo
    deck.panel(slide, 0, 0, g.width, 0.16, C_ACCENT, shape=MSO_SHAPE.RECTANGLE)
    deck.text(slide, "Perturb-seq QC report", g.margin, 2.02,
              g.content_width, 0.28, size=13, bold=True, color=C_ACCENT,
              caps=True)
    txt, size = fit_text(title, g.content_width, 1.5, 40, 22)
    deck.text(slide, txt, g.margin, 2.45, g.content_width, 1.6,
              size=size, bold=True, color=C_INK)
    if subtitle:
        sub, ssize = fit_text(subtitle, g.content_width * 0.8, 0.8, 15, MIN_PT)
        deck.text(slide, sub, g.margin, 4.15, g.content_width * 0.8, 0.85,
                  size=ssize, color=C_INK_SOFT)
    deck.rule(slide, g.margin, 5.25, g.content_width)
    deck.text(slide, "   ·   ".join(meta), g.margin, 5.42,
              g.content_width, 0.7, size=MIN_PT, color=C_INK_FAINT)
    return slide


def contents_slide(deck: Deck, entries: Sequence[tuple[int, str, str]]):
    slide = deck.new("Contents")
    g = deck.geo
    n = len(entries)
    cols = 2 if n > 5 else 1
    per = math.ceil(n / cols)
    col_w = (g.content_width - 0.45) / cols
    row_h = min(1.05, g.content_height / per)
    for c in range(cols):
        chunk = entries[c * per:(c + 1) * per]
        y = g.content_top
        x = g.content_left + c * (col_w + 0.45)
        for num, title, blurb in chunk:
            head, hsize = fit_text(f"{num}.  {title}", col_w, 0.3, 14, MIN_PT)
            deck.text(slide, head, x, y, col_w, 0.3, size=hsize, bold=True,
                      color=C_INK)
            body, size = fit_text(blurb, col_w - 0.32, row_h - 0.42, MIN_PT,
                                  MIN_PT)
            deck.text(slide, body, x + 0.32, y + 0.32, col_w - 0.32,
                      row_h - 0.4, size=size, color=C_INK_FAINT)
            y += row_h
    return slide


def _metric_grid(n: int, width: float, height: float | None = None
                 ) -> tuple[int, int, float, float, float]:
    """Grid geometry for n metric cards.

    Cards are as wide as they can be while still fitting the height available:
    a run with 13 headline numbers gets four roomy columns, one with 17 gets
    five slightly narrower ones, and neither spills a lonely card onto a second
    slide.
    """
    # Tall enough for a value plus three lines of 12pt label: "Cells with an
    # assigned guide" wraps, and a label that wraps past the card edge is worse
    # than a narrower card.
    gap, card_h = 0.14, 1.25
    widest = max(2, int(width // 2.7))
    per_row = min(n, widest)
    if height:
        while per_row < n:
            rows = math.ceil(n / per_row)
            if rows * (card_h + gap) <= height:
                break
            per_row += 1
        per_row = min(per_row, max(2, int(width // 1.9)))
    rows = math.ceil(n / per_row)
    card_w = (width - gap * (per_row - 1)) / per_row
    return per_row, rows, card_w, card_h, gap


def metric_cards(deck: Deck, slide, arts: Sequence[Art], left: float,
                 top: float, width: float, height: float | None = None) -> float:
    """A row (or rows) of value/label cards, as in the report's .metrics grid."""
    if not arts:
        return 0.0
    per_row, rows, card_w, card_h, gap = _metric_grid(len(arts), width, height)
    for i, a in enumerate(arts):
        r, c = divmod(i, per_row)
        x = left + c * (card_w + gap)
        y = top + r * (card_h + gap)
        level = str(a.data.get("level", "info"))
        ink, soft = LEVELS.get(level, LEVELS["info"])
        deck.panel(slide, x, y, card_w, card_h, soft, ink)
        value, _ = fmt_value(a.data.get("value"))
        unit = plain(a.data.get("unit", ""))
        shown = f"{value}{unit}" if unit else value
        vtxt, vsize = fit_text(shown, card_w - 0.28, 0.44, 24, 14)
        deck.text(slide, vtxt, x + 0.14, y + 0.10, card_w - 0.28, 0.46,
                  size=vsize, bold=True, color=ink)
        ltxt, lsize = fit_text(plain(a.title), card_w - 0.28, 0.64, MIN_PT,
                               MIN_PT)
        deck.text(slide, ltxt, x + 0.14, y + 0.56, card_w - 0.28, 0.64,
                  size=lsize, color=C_INK_SOFT)
    return rows * (card_h + gap)


def _note_layout(body: str, width: float, max_height: float,
                 max_lines: int = 4) -> tuple[str, float, float, bool]:
    """(text shown, size, panel height, whether anything was left out)."""
    text = plain(body)
    size = MIN_PT
    budget = min(max_height - 0.58, max_lines * size * _LINE / 72.0)
    body_txt, trimmed = fit_sentences(text, width - 0.44, budget, size)
    lines = _line_count(body_txt, width - 0.44, size)
    # A spare line for long blocks only: LibreOffice and Google Slides wrap a
    # little differently from the estimate, and a note that overflows its panel
    # looks broken -- but padding every two-line note costs a whole slide.
    slack = 1 if lines >= 3 else 0
    h = 0.34 + (lines + slack) * size * _LINE / 72.0 + 0.14
    return body_txt, size, h, trimmed


def note_block(deck: Deck, slide, title: str, body: str, level: str,
               left: float, top: float, width: float,
               max_height: float, max_lines: int = 4) -> float:
    ink, soft = LEVELS.get(level, LEVELS["info"])
    body_txt, size, h, trimmed = _note_layout(body, width, max_height, max_lines)
    if trimmed:
        # Whatever did not fit is still readable, just not on the slide.
        existing = (slide.notes_slide.notes_text_frame.text
                    if slide.has_notes_slide else "")
        deck.notes(slide, (existing + "\n\n" if existing else "")
                   + f"{plain(title)}\n{plain(body)}")
    deck.panel(slide, left, top, width, h, soft, ink)
    deck.text(slide, plain(title), left + 0.22, top + 0.07, width - 0.44, 0.24,
              size=MIN_PT, bold=True, color=ink, caps=True)
    deck.text(slide, body_txt, left + 0.22, top + 0.33, width - 0.44,
              h - 0.42, size=size, color=C_INK_SOFT)
    return h


# ===========================================================================
# Section prose: blurb, metric cards and method notes as a flowable block
# ===========================================================================
Block = tuple[str, Any]


def build_block(blurb: str, metrics: Sequence[Art], notes: Sequence[Art],
                skipped: Sequence[Art]) -> list[Block]:
    """The non-figure content of a section, in the order the HTML shows it."""
    items: list[Block] = []
    if plain(blurb):
        items.append(("blurb", plain(blurb)))
    if metrics:
        items.append(("metrics", list(metrics)))
    for a in sorted(notes, key=lambda x: x.order):
        body = a.data.get("body") or a.caption
        items.append(("note", (a.title, body, str(a.data.get("level", "info")))))
    if skipped:
        # "Registered, not produced, and here is why" is itself a QC finding, so
        # it travels with the section rather than being silently dropped.
        # First sentence only: the reason a figure is absent is usually one
        # clause followed by paragraphs of rationale.
        body = "\n".join(
            f"•  {plain(a.title)}: {(sentences(a.skipped_reason) or [''])[0]}"
            for a in skipped
        )
        items.append(("skipped", ("Registered but not produced", body, "info")))
    return items


def _blurb_layout(text: str, width: float) -> tuple[str, float, float]:
    txt, size = fit_text(text, width, 0.62, 13, MIN_PT)
    return txt, size, 0.18 + _line_count(txt, width, size) * size * _LINE / 72.0


def measure_block(items: Sequence[Block], width: float,
                  note_cap: float = 2.4, max_lines: int = 4) -> float:
    total = 0.0
    for kind, payload in items:
        if kind == "blurb":
            total += _blurb_layout(payload, width)[2]
        elif kind == "metrics":
            _, rows, _, card_h, gap = _metric_grid(len(payload), width)
            total += rows * (card_h + gap) + 0.16
        else:
            total += _note_layout(payload[1], width, note_cap, max_lines)[2] + 0.14
    return total


def draw_block(deck: Deck, slide, items: Sequence[Block], left: float,
               top: float, width: float, max_h: float,
               max_lines: int = 4) -> tuple[float, list[Block]]:
    """Draw as much of the block as fits; return (height used, what is left)."""
    y, bottom = top, top + max_h
    for i, (kind, payload) in enumerate(items):
        if kind == "blurb":
            txt, size, h = _blurb_layout(payload, width)
            if h > bottom - y:
                return y - top, list(items[i:])
            deck.text(slide, txt, left, y, width, h, size=size, italic=True,
                      color=C_INK_SOFT)
            y += h
        elif kind == "metrics":
            per_row, rows, _, card_h, gap = _metric_grid(
                len(payload), width, bottom - y - 0.16)
            h = rows * (card_h + gap) + 0.16
            if h > bottom - y:
                fits = int((bottom - y - 0.16) // (card_h + gap))
                if fits < 1:
                    return y - top, list(items[i:])
                here = payload[:fits * per_row]
                rest = payload[fits * per_row:]
                y += metric_cards(deck, slide, here, left, y, width,
                                  bottom - y - 0.16) + 0.16
                remaining: list[Block] = ([("metrics", rest)] if rest else [])
                return y - top, remaining + list(items[i + 1:])
            y += metric_cards(deck, slide, payload, left, y, width,
                              bottom - y - 0.16) + 0.16
        else:
            title, body, level = payload
            if bottom - y < 0.75:  # too little room to be legible: overleaf
                return y - top, list(items[i:])
            y += note_block(deck, slide, title, body, level, left, y, width,
                            bottom - y, max_lines) + 0.14
    return y - top, []


def section_slides(deck: Deck, title: str, num: int,
                   items: Sequence[Block]) -> list:
    """Standalone section slide(s) for a section's prose, metrics and notes."""
    g = deck.geo
    # Prefer four lines a note, but tighten to three or two before letting a
    # section's prose run onto a second slide -- the full text is in the notes.
    max_lines = 4
    for cap in (4, 3, 2):
        if measure_block(items, g.content_width, max_lines=cap) <= g.content_height:
            max_lines = cap
            break
        max_lines = cap
    slides, pending, first = [], list(items), True
    while pending:
        slide = deck.new(f"{num}. {title}" + ("" if first else " (cont.)"))
        slides.append(slide)
        used, pending = draw_block(deck, slide, pending, g.content_left,
                                   g.content_top, g.content_width,
                                   g.content_height, max_lines)
        if used <= 0:  # nothing fit at all; bail rather than loop forever
            break
        first = False
    return slides


# ===========================================================================
# Figure grouping
# ===========================================================================
@dataclass
class FigGroup:
    """The figures that share one slide, and how to label them."""

    items: list[Art]
    family: str = ""             # common prefix, shown above the grid
    labels: list[str] = field(default_factory=list)   # per-cell headings
    shared_caption: str = ""     # one caption for the whole grid

    def __post_init__(self):
        if not self.labels:
            self.labels = [plain(a.title) for a in self.items]


_BY = re.compile(r"^(.*?) by (.+)$")


def _family(a: Art) -> tuple[str, str] | None:
    """Split "Guide purity by fixation" into ("Guide purity", "fixation")."""
    m = _BY.match(plain(a.title))
    return (m.group(1).strip(), m.group(2).strip()) if m else None


def _similar_size(a: tuple[int, int], b: tuple[int, int], tol: float = 0.08) -> bool:
    return all(
        abs(x - y) <= tol * max(x, y, 1) for x, y in zip(a, b)
    )


def group_figures(figs: Sequence[Art], per_slide: int, root: Path,
                  pack: str = "family", pack_aspect: float = 1.7,
                  family_max: int = 6) -> list[FigGroup]:
    """Decide which figures share a slide.

    Three strategies:

    ``width``   mirror the HTML exactly -- ``full`` alone, ``half`` in pairs.
    ``aspect``  ignore the declared width and pair anything that is not
                genuinely wide, judged from the PNG's real aspect ratio.
    ``family``  ``aspect``, plus: a run of "<something> by <condition>" figures
                that are the same plot at the same size becomes one grid slide
                labelled by condition. This is the clustering the report implies
                but cannot express in a single column of HTML, and it is where
                the slide count actually comes down: five such families in a
                typical run turn twenty slides into five.
    """
    def size(a: Art) -> tuple[int, int]:
        p = a.resolve(root)
        return image_size(p) if (p and p.exists()) else (1600, 1200)

    def is_full(a: Art) -> bool:
        if pack == "width":
            return a.width == "full"
        w, h = size(a)
        return (w / max(h, 1)) > pack_aspect

    groups: list[FigGroup] = []
    i, n = 0, len(figs)
    while i < n:
        a = figs[i]
        fam = _family(a) if pack == "family" else None
        if fam and family_max > 1:
            # Greedily extend the family: same prefix, same section, same size.
            run, labels = [a], [fam[1]]
            j = i + 1
            while j < n and len(run) < family_max:
                nxt, nfam = figs[j], _family(figs[j])
                if (nfam and nfam[0] == fam[0] and nxt.section == a.section
                        and _similar_size(size(nxt), size(a))):
                    run.append(nxt)
                    labels.append(nfam[1])
                    j += 1
                else:
                    break
            if len(run) >= 3:
                caps = {plain(x.caption) for x in run}
                groups.append(FigGroup(
                    items=run, family=fam[0], labels=labels,
                    shared_caption=caps.pop() if len(caps) == 1 else "",
                ))
                i = j
                continue

        # Not a family: fall back to width/aspect packing.
        if is_full(a) or per_slide == 1:
            groups.append(FigGroup(items=[a]))
            i += 1
            continue
        run = [a]
        j = i + 1
        while j < n and len(run) < per_slide and not is_full(figs[j]) \
                and not (_family(figs[j]) and pack == "family"):
            run.append(figs[j])
            j += 1
        if len(run) == 3 and any(len(plain(x.caption)) > 200 for x in run):
            run = run[:2]
        groups.append(FigGroup(items=run))
        i += len(run)
    return groups


def best_grid(n: int, area_w: float, area_h: float, aspect: float,
              gap: float, title_h: float) -> tuple[int, int]:
    """Rows and columns that render the figures as large as possible."""
    best = (0.0, 1, n)
    for cols in range(1, n + 1):
        rows = math.ceil(n / cols)
        cw = (area_w - gap * (cols - 1)) / cols
        ch = (area_h - gap * (rows - 1)) / rows - title_h
        if cw <= 0.8 or ch <= 0.45:
            continue
        drawn_h = min(cw / aspect, ch)
        score = drawn_h * drawn_h * aspect * n
        if score > best[0]:
            best = (score, cols, rows)
    return best[1], best[2]


def refine_groups(groups: Sequence[FigGroup], root: Path, geo: Geometry,
                  min_panel_w: float = 3.6) -> list[FigGroup]:
    """Split any group whose panels would render too small to read.

    Packing four square panels onto one slide is a false economy: the grid that
    maximises area still leaves each one about two inches wide. Where that
    happens the group is halved, which costs a slide and returns the panels to a
    size someone can actually look at.
    """
    out: list[FigGroup] = []
    for grp in groups:
        items = grp.items
        if len(items) < 3:
            out.append(grp)
            continue
        sizes = []
        for a in items:
            path = a.resolve(root)
            sizes.append(image_size(path) if (path and path.exists()) else (1600, 1200))
        aspect = sum(w / max(h, 1) for w, h in sizes) / len(sizes)
        # Height a grid slide has after the standfirst and the family label.
        grid_h = geo.content_height - 0.42 - (0.36 if grp.family else 0.0)
        cols, rows = best_grid(len(items), geo.content_width, grid_h, aspect,
                               0.24, 0.32)
        cell_w = (geo.content_width - 0.24 * (cols - 1)) / cols
        cell_h = (grid_h - 0.24 * (rows - 1)) / rows - 0.32
        drawn_w = min(cell_w, cell_h * aspect)
        if drawn_w >= min_panel_w:
            out.append(grp)
            continue
        half = math.ceil(len(items) / 2)
        for start in range(0, len(items), half):
            chunk = list(range(start, min(start + half, len(items))))
            out.append(FigGroup(
                items=[items[i] for i in chunk],
                family=grp.family,
                labels=[grp.labels[i] for i in chunk],
                shared_caption=grp.shared_caption,
            ))
    return out


def figure_slides(deck: Deck, title: str, num: int,
                  groups: Sequence[FigGroup], root: Path, cap_dir: Path,
                  max_px: int, cont_start: bool,
                  banner: Sequence[Block] = ()):
    """Figure slides for one section, laid out as a grid.

    ``banner`` is prose/metrics small enough to share the first figure slide --
    a section whose only non-figure content is a one-line blurb and a single
    metric does not deserve a slide of its own.
    """
    g = deck.geo
    for gi, group in enumerate(refine_groups(groups, root, g)):
        items = group.items
        slide = deck.new(f"{num}. {title}"
                         + (" (cont.)" if (cont_start or gi) else ""))
        top, avail_h = g.content_top, g.content_height

        # Standfirst first: it is anchored under the title, so anything else on
        # the slide has to start below whatever height it consumed. (Drawing the
        # banner first put the section blurb and the standfirst in the same
        # 12-inch strip.)
        shared_cap = plain(group.shared_caption) or (
            plain(items[0].caption) if len(items) == 1 else ""
        )
        if shared_cap:
            shows, verdict, label = caption_brief(shared_cap)
        elif group.family:
            # Family whose captions differ only in the level they name: say what
            # the grid is, then quote only the guidance all of them share.
            shows = (f"{group.family} for each level of "
                     f"{', '.join(group.labels)}.")
            _, verdict, label = caption_brief(
                shows + " " + common_sentences([a.caption for a in items])
            )
        else:
            shows, verdict, label = (
                " · ".join(plain(a.title) for a in items), "", "",
            )
        used_brief = deck.brief(slide, shows, verdict, label)
        top += used_brief
        avail_h -= used_brief

        if gi == 0 and banner:
            used, _ = draw_block(deck, slide, banner, g.content_left, top,
                                 g.content_width, avail_h * 0.42)
            top += used + 0.14
            avail_h -= used + 0.14

        n = len(items)
        if group.family:
            fam_txt, fam_size = fit_text(group.family, g.content_width, 0.32,
                                         14, MIN_PT)
            deck.text(slide, fam_txt, g.content_left, top, g.content_width, 0.32,
                      size=fam_size, bold=True, color=C_INK)
            top += 0.36
            avail_h -= 0.36

        # With the standfirst carrying the interpretation, a caption is repeated
        # below a panel only when the panels differ and each needs its own.
        shared, shared_h = "", 0.0
        per_cell_caps = 0 if (shared_cap or n >= 4) else {1: 3, 2: 3, 3: 2}[n]

        gap = 0.24
        title_h = 0.32
        sizes = []
        for a in items:
            path = a.resolve(root)
            sizes.append(image_size(path) if (path and path.exists()) else (1600, 1200))
        aspect = sum(w / max(h, 1) for w, h in sizes) / len(sizes)
        grid_h = avail_h - shared_h
        cols, rows = best_grid(n, g.content_width, grid_h, aspect, gap, title_h)
        cell_w = (g.content_width - gap * (cols - 1)) / cols
        cell_h = (grid_h - gap * (rows - 1)) / rows

        notes_out = []
        tallest_row_bottom = top
        for idx, a in enumerate(items):
            r, c = divmod(idx, cols)
            x = g.content_left + c * (cell_w + gap)
            y = top + r * (cell_h + gap)
            label = group.labels[idx] if idx < len(group.labels) else plain(a.title)
            ltxt, lsize = fit_text(label, cell_w, title_h, 14 if n <= 2 else 13,
                                   MIN_PT)
            deck.text(slide, ltxt, x, y, cell_w, title_h, size=lsize, bold=True,
                      color=C_INK)
            y += title_h
            room_h = cell_h - title_h
            path = a.resolve(root)
            full_caption = plain(a.caption)
            if path is None or not path.exists():
                reason = plain(a.skipped_reason) or (
                    f"registered {a.path} but the file is not on disk"
                )
                box_h = min(1.1, room_h)
                deck.panel(slide, x, y, cell_w, box_h, C_LINE_SOFT, C_LINE)
                deck.text(slide, f"{plain(a.title)} — not produced\n{reason}",
                          x + 0.18, y + 0.18, cell_w - 0.36, box_h - 0.36,
                          size=MIN_PT, color=C_INK_FAINT)
                notes_out.append(f"{plain(a.title)}: not produced. {reason}")
                continue
            cap_h = 0.0
            if per_cell_caps:
                cap_h = 0.62 if n == 1 else 0.86
            src = maybe_downscale(path, max_px, cap_dir)
            used = deck.picture(slide, src, x, y, cell_w, room_h - cap_h)
            if per_cell_caps:
                cap_top = y + used + 0.12
                box = max(cap_h, (y + room_h) - cap_top)
                ctxt, _cut = fit_sentences(full_caption, cell_w, box, MIN_PT)
                if _line_count(ctxt, cell_w, MIN_PT) * MIN_PT * _LINE / 72.0 > box:
                    ctxt = ""      # not even one sentence fits: notes carry it
                if ctxt:
                    deck.text(slide, ctxt, x, cap_top, cell_w, box, size=MIN_PT,
                              color=C_INK_SOFT)
            tallest_row_bottom = max(tallest_row_bottom, y + used)
            if full_caption:
                notes_out.append(f"{plain(a.title)}\n{full_caption}")

        if shared:
            cap_top = min(tallest_row_bottom + 0.14,
                          top + grid_h + (0.0 if rows == 1 else 0.0))
            box = (g.content_top + g.content_height) - cap_top
            ctxt, csize = fit_text(shared, g.content_width, max(0.3, box), 10, 7.5)
            deck.text(slide, ctxt, g.content_left, cap_top, g.content_width,
                      max(0.3, box), size=csize, color=C_INK_SOFT)
        deck.notes(slide, "\n\n".join(notes_out))


def table_rows(a: Art, root: Path, max_rows: int) -> tuple[list[str], list[dict], int]:
    """Inline rows if the registry carried them, else read the CSV it points at."""
    rows = list(a.data.get("rows") or [])
    cols = list(a.data.get("columns") or [])
    total = len(rows)
    if not rows:
        path = a.resolve(root)
        if path and path.exists() and path.suffix.lower() in (".csv", ".tsv"):
            delim = "\t" if path.suffix.lower() == ".tsv" else ","
            with path.open(newline="", encoding="utf-8") as fh:
                reader = csv.DictReader(fh, delimiter=delim)
                cols = cols or list(reader.fieldnames or [])
                for i, r in enumerate(reader):
                    total = i + 1
                    if i < max_rows * 4:  # bounded read; we only show max_rows
                        rows.append(r)
    if not cols and rows:
        cols = list(rows[0].keys())
    return cols, rows, total


def table_slides(deck: Deck, title: str, num: int, tables: Sequence[Art],
                 root: Path, max_rows: int, cont: bool, max_cols: int = 8):
    """Table slides, paginated over rows *and* columns.

    The sequencing-metrics table has ~25 columns; squeezing those into 12 inches
    gives you a slide nobody can read. Wide tables are split into column groups
    that repeat the first column as the row key, capped at three groups and
    three row pages -- past that the CSV is the right medium and the footnote
    says so.
    """
    g = deck.geo
    for a in tables:
        cols, rows, total = table_rows(a, root, max_rows)
        if not cols:
            continue
        if len(cols) <= max_cols:
            col_pages = [list(cols)]
        else:
            key, rest = cols[0], list(cols[1:])
            step = max(1, max_cols - 1)
            col_pages = [
                [key] + rest[i:i + step] for i in range(0, len(rest), step)
            ]
        col_pages, dropped_cols = col_pages[:3], max(0, len(col_pages[3:]))
        row_pages = [
            rows[i:i + max_rows]
            for i in range(0, min(len(rows), max_rows * 3), max_rows)
        ] or [[]]
        shown_rows = sum(len(pg) for pg in row_pages)
        shown_cols = len({c for pg in col_pages for c in pg})

        for ci, page_cols in enumerate(col_pages):
            for ri, page in enumerate(row_pages):
                slide = deck.new(f"{num}. {title}"
                                 + (" (cont.)" if (cont or ci or ri) else ""))
                cont = True
                y = g.content_top
                shows, verdict, vlabel = caption_brief(plain(a.caption))
                if not shows:
                    shows = (f"{plain(a.title)}: {total:,} rows × "
                             f"{len(cols)} columns")
                y += deck.brief(slide, shows, verdict, vlabel)
                label = plain(a.title)
                if len(row_pages) > 1:
                    label += f"  (rows {ri + 1}/{len(row_pages)})"
                if len(col_pages) > 1:
                    label += f"  (cols {ci + 1}/{len(col_pages)})"
                ltxt, lsize = fit_text(label, g.content_width, 0.32, 14,
                                       MIN_PT)
                deck.text(slide, ltxt, g.content_left, y, g.content_width, 0.32,
                          size=lsize, bold=True, color=C_INK)
                y += 0.36
                foot = []
                if total > shown_rows:
                    foot.append(f"{shown_rows:,} of {total:,} rows")
                if len(cols) > shown_cols:
                    foot.append(f"{shown_cols} of {len(cols)} columns")
                _draw_table(deck, slide, page_cols, page, g.content_left, y,
                            g.content_width,
                            (g.content_top + g.content_height) - y
                            - (0.34 if foot else 0.0))
                if foot:
                    deck.text(
                        slide,
                        "Showing " + " and ".join(foot)
                        + "; the full table is in the CSV alongside the report.",
                        g.content_left, g.content_top + g.content_height - 0.26,
                        g.content_width, 0.26, size=MIN_PT, italic=True,
                        color=C_INK_FAINT,
                    )
                deck.notes(slide, plain(a.caption))


def _draw_table(deck: Deck, slide, cols: Sequence[str], rows: Sequence[dict],
                left: float, top: float, width: float, height: float):
    n_rows = len(rows) + 1
    # 12pt needs ~0.34in of row; the header gets two lines for long column names.
    head_h = 0.52
    row_h = min(0.40, max(0.34, (height - head_h) / max(1, len(rows))))
    shape = slide.shapes.add_table(
        n_rows, len(cols), Inches(left), Inches(top), Inches(width),
        Inches(head_h + row_h * len(rows)),
    )
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = True
    # Column widths proportional to the widest cell each column has to hold.
    weights = []
    for c in cols:
        w = len(str(c))
        for r in rows:
            w = max(w, len(fmt_value(r.get(c))[0]))
        weights.append(max(4, min(w, 34)))
    tot = sum(weights)
    for i, wgt in enumerate(weights):
        tbl.columns[i].width = Emu(int(Inches(width) * wgt / tot))
    tbl.rows[0].height = Inches(head_h)
    for i in range(len(rows)):
        tbl.rows[i + 1].height = Inches(row_h)

    def cell_text(cell, txt: str, *, bold=False, color=C_INK, align=PP_ALIGN.LEFT,
                  size=MIN_PT):
        cell.margin_left = cell.margin_right = Inches(0.06)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    for j, c in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_ACCENT_SOFT
        cell_text(cell, plain(c), bold=True, color=C_ACCENT)
    # One line per cell at 12pt: how many characters that is depends on the
    # column's own width, so each column gets its own budget.
    budgets = [
        max(6, int(Emu(tbl.columns[j].width).inches * 72 / (_CHAR_EM * MIN_PT)))
        for j in range(len(cols))
    ]
    for i, r in enumerate(rows):
        for j, c in enumerate(cols):
            txt, numeric = fmt_value(r.get(c))
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE if i % 2 == 0 else C_LINE_SOFT
            cell_text(cell, truncate(txt, budgets[j]), color=C_INK_SOFT,
                      align=PP_ALIGN.RIGHT if numeric else PP_ALIGN.LEFT)


def provenance_slides(deck: Deck, num: int, rows: Sequence[dict],
                      per_slide: int = 12):
    g = deck.geo
    pages = [rows[i:i + per_slide] for i in range(0, len(rows), per_slide)]
    for pi, page in enumerate(pages):
        slide = deck.new(f"{num}. Run provenance"
                         + (f"  ({pi + 1}/{len(pages)})" if len(pages) > 1 else ""))
        _draw_table(deck, slide, ["item", "value"], page, g.content_left,
                    g.content_top, g.content_width, g.content_height)


def conclusions_slide(deck: Deck, headings: Sequence[str]):
    """A deliberately empty closing slide, for the reader to fill in.

    The report can say whether the assay worked; only a person can say what to
    do about it. Real text boxes with a faint prompt, so they can be typed into
    directly in Google Slides and deleted if unused.
    """
    g = deck.geo
    slide = deck.new("Conclusions & next steps")
    n = len(headings)
    gap = 0.22
    h = (g.content_height - gap * (n - 1)) / n
    prompts = [
        "Click to add what the data supports.",
        "Click to add what is still uncertain, and what would settle it.",
        "Click to add owners and dates.",
    ]
    for i, head in enumerate(headings):
        y = g.content_top + i * (h + gap)
        deck.panel(slide, g.content_left, y, g.content_width, h,
                   C_WHITE, C_LINE)
        deck.text(slide, plain(head), g.content_left + 0.24, y + 0.16,
                  g.content_width - 0.48, 0.28, size=13, bold=True,
                  color=C_ACCENT, caps=True)
        deck.text(slide, prompts[i] if i < len(prompts) else "Click to add text.",
                  g.content_left + 0.24, y + 0.52, g.content_width - 0.48,
                  h - 0.68, size=MIN_PT, italic=True, color=C_INK_FAINT)
    deck.notes(slide, "Placeholder slide: replace the prompt text with the "
                      "conclusions and next steps agreed in the meeting.")
    return slide


def checklist_slide(deck: Deck, num: int):
    """The report's 'before treating these results as final' list, if available."""
    try:
        from perturbseq_report import text as T  # optional
        items = [plain(x) for x in getattr(T, "FINAL_CHECKLIST", [])]
    except Exception:
        return None
    if not items:
        return None
    g = deck.geo
    slide = deck.new(f"{num}. Before treating these results as final")
    y = g.content_top
    deck.text(
        slide,
        "This report can tell you whether the assay worked. It cannot tell you "
        "whether the conclusion you want to draw from it is sound. These are the "
        "checks it cannot do for you.",
        g.content_left, y, g.content_width, 0.6, size=MIN_PT, italic=True,
        color=C_INK_SOFT,
    )
    y += 0.62
    room = (g.content_top + g.content_height) - y
    per = room / max(1, len(items))
    text = "\n".join(f"•  {i}" for i in items)
    txt, size = fit_text(text, g.content_width, room, 13, MIN_PT)
    deck.text(slide, txt, g.content_left, y, g.content_width, room, size=size,
              color=C_INK, spacing=4)
    return slide


# ===========================================================================
# Orchestration
# ===========================================================================
def build_deck(
    artifacts_json: Path,
    out_path: Path,
    title: str | None = None,
    subtitle: str = "",
    aspect: str = "16:9",
    per_slide: int = 2,
    table_mode: str = "charts",
    max_table_rows: int = 11,
    max_table_cols: int = 5,
    drop_redundant_tables: bool = True,
    no_table_sections: Sequence[str] = ("perturbation",),
    essential_notes_only: bool = True,
    max_notes_per_section: int = 4,
    max_charts_per_section: int = 3,
    drop_tables: Sequence[str] = (),
    keep_tables: Sequence[str] = (),
    conclusions: Sequence[str] = (
        "Conclusions", "Open questions", "Next steps",
    ),
    max_image_px: int = 1800,
    banner_max: float = 1.45,
    pack: str = "family",
    pack_aspect: float = 1.7,
    family_max: int = 6,
    drop_duplicates: bool = True,
    provenance: Sequence[dict] = (),
    only: Sequence[str] | None = None,
    appendix: bool = False,
    skip: Sequence[str] | None = None,
    contents: bool = True,
    quiet: bool = False,
) -> Path:
    reg = Registry.load(artifacts_json)
    by_section = reg.by_section()
    titles = {k: t for k, t, _ in reg.sections}
    blurbs = {k: b for k, _, b in reg.sections}
    spine = [k for k, _, _ in reg.sections]

    wanted = [k for k in spine if k in by_section]
    if only:
        wanted = [k for k in wanted if k in set(only)]
    if skip:
        wanted = [k for k in wanted if k not in set(skip)]
    if not appendix:
        # Configuration, provenance and input-file bookkeeping: archived in the
        # HTML report, not something a meeting reads off a slide.
        wanted = [k for k in wanted if k != "appendix"]
    # A section with nothing produced and nothing explained is not worth a slide.
    wanted = [
        k for k in wanted
        if any(a.produced for a in by_section[k])
        or any(a.skipped_reason for a in by_section[k])
    ]

    run_label = title or f"Perturb-seq QC — {reg.root.resolve().name}"
    generated = datetime.now().strftime("%Y-%m-%d %H:%M")
    deck = Deck(aspect, truncate(run_label, 72), generated)

    n_fig = sum(
        1 for a in reg.artifacts
        if a.kind == "figure" and a.produced
        and (a.resolve(reg.root) or Path("/nonexistent")).exists()
    )
    n_skip = sum(1 for a in reg.artifacts if not a.produced)
    meta = [f"generated {generated}", f"{n_fig} figures",
            f"{len(wanted)} sections"]
    if n_skip:
        meta.append(f"{n_skip} not applicable")
    prov = {r["item"]: r["value"] for r in provenance}
    for k in ("pipeline version", "cells analysed", "backend"):
        if prov.get(k):
            meta.append(f"{k}: {prov[k]}")
    title_slide(
        deck, run_label,
        subtitle or ("Figures, metrics and captions carried over from the HTML QC "
                     "report. Import into Google Slides with File › Import "
                     "slides."),
        meta,
    )

    numbering = {k: i + 1 for i, k in enumerate(wanted)}
    if contents and len(wanted) > 2:
        contents_slide(
            deck, [(numbering[k], titles.get(k, k), plain(blurbs.get(k, "")))
                   for k in wanted]
        )

    seen_digests: dict[str, str] = {}
    duplicates: list[tuple[str, str]] = []
    dropped_tables: list[tuple[str, str]] = []
    dropped_notes: list[tuple[str, str]] = []
    charted_tables: list[tuple[str, str]] = []
    unchartable: list[tuple[str, str]] = []

    with tempfile.TemporaryDirectory(prefix="ps_slides_") as tmp:
        cache = Path(tmp)
        for key in wanted:
            items = by_section[key]
            num = numbering[key]
            name = titles.get(key, key)
            metrics = [a for a in items if a.kind == "metric" and a.produced]
            notes = [a for a in items if a.kind == "note" and a.produced]
            if essential_notes_only:
                notes = curate_notes(notes, max_notes_per_section, dropped_notes)
            figs = [a for a in items if a.kind == "figure" and a.produced]
            if drop_duplicates:
                kept = []
                for a in figs:
                    path = a.resolve(reg.root)
                    dig = file_digest(path) if path and path.exists() else None
                    if dig and dig in seen_digests:
                        # Byte-identical to a figure already in the deck: the
                        # old pipeline wrote some plots twice, and a duplicate
                        # image carries no information a second slide could add.
                        duplicates.append((a.title, seen_digests[dig]))
                        continue
                    if dig:
                        seen_digests[dig] = plain(a.title)
                    kept.append(a)
                figs = kept
            tables = [a for a in items if a.kind == "table" and a.produced]
            if tables and table_mode == "full" and key in set(no_table_sections):
                for t in tables:
                    dropped_tables.append(
                        (plain(t.title), f"--no-tables-in {key}")
                    )
                tables = []
            if tables:
                fig_titles = [a.title for a in items if a.kind == "figure"]
                keep = []
                for t in tables:
                    title = plain(t.title)
                    if any(re.search(p, title, re.I) for p in keep_tables):
                        keep.append(t)
                        continue
                    why = None
                    for pat in drop_tables:
                        if re.search(pat, title, re.I):
                            why = f"matches --drop-tables {pat!r}"
                            break
                    if why is None and drop_redundant_tables:
                        for pat, reason in REDUNDANT_TABLES:
                            if re.search(pat, title, re.I):
                                why = reason
                                break
                    if why is None and drop_redundant_tables:
                        shown_in = table_is_plotted(t.title, fig_titles)
                        if shown_in:
                            why = f"plotted in {shown_in!r}"
                    if why:
                        dropped_tables.append((title, why))
                    else:
                        keep.append(t)
                tables = keep
            skipped = [a for a in items if not a.produced]

            block = build_block(blurbs.get(key, ""), metrics, notes, skipped)
            if table_mode == "charts" and tables:
                note_titles = {plain(a.title).lower() for a in notes}
                chartable = []
                for t in tables:
                    if plain(t.title).lower() in note_titles:
                        unchartable.append(
                            (plain(t.title),
                             "a method note of the same name already reports it")
                        )
                    else:
                        chartable.append(t)
                new_figs = tables_as_figures(
                    chartable[:max_charts_per_section] if max_charts_per_section
                    else chartable,
                    reg.root, cache / "charts", max_table_rows,
                    charted_tables, unchartable,
                )
                for extra in chartable[max_charts_per_section:] if max_charts_per_section else []:
                    unchartable.append(
                        (plain(extra.title),
                         f"over the {max_charts_per_section}-chart budget for "
                         f"its section")
                    )
                figs = figs + new_figs
            groups = group_figures(
                figs, per_slide, reg.root, pack, pack_aspect, family_max
            ) if figs else []
            # A thin block (a blurb plus a metric or two) rides on the first
            # figure slide instead of claiming a mostly-empty slide of its own.
            block_h = measure_block(block, deck.geo.content_width) if block else 0.0
            banner = (
                bool(block) and bool(groups)
                and block_h <= banner_max
                # 0.5in for the standfirst, 0.14 gap: what is left has to still
                # be worth showing a figure in.
                and (deck.geo.content_height - 0.5 - block_h - 0.14) >= 3.2
            )
            opened = bool(block) and not banner
            if opened:
                section_slides(deck, name, num, block)
            if groups:
                figure_slides(deck, name, num, groups, reg.root, cache,
                              max_image_px, cont_start=opened,
                              banner=block if banner else ())
            if table_mode == "full" and tables:
                table_slides(deck, name, num, tables, reg.root, max_table_rows,
                             cont=bool(opened or groups),
                             max_cols=max_table_cols)

        appendix_num = len(wanted) + 1
        if appendix:
            checklist_slide(deck, appendix_num)
        already = any(
            a.kind == "table"
            and [str(c).lower() for c in (a.data.get("columns") or [])]
            == ["item", "value"]
            for a in reg.artifacts
        )
        if appendix and provenance and not already:
            provenance_slides(deck, appendix_num, provenance)
        if conclusions:
            conclusions_slide(deck, conclusions)
        out = deck.save(out_path)

    if not quiet:
        charts = f" + {len(charted_tables)} charted tables" if charted_tables else ""
        print(f"{out}  —  {deck.n} slides, "
              f"{n_fig - len(duplicates)} figures{charts}")
        for title, first in duplicates:
            print(f"  dropped figure '{plain(title)}': byte-identical to '{first}'")
        for title, why in dropped_tables:
            print(f"  dropped table  '{title}': {why}")
        for title, why in dropped_notes:
            print(f"  dropped note   '{title}': {why}")
        for title, what in charted_tables:
            print(f"  charted table  '{title}': {what}")
        for title, why in unchartable:
            print(f"  dropped table  '{title}': {why}")
        print("Open in Google Slides: upload to Drive and 'Open with > Google "
              "Slides', or File > Import slides in an existing deck.")
    return out


def resolve_input(arg: Path) -> Path:
    if arg.is_dir():
        cand = arg / "artifacts.json"
        if not cand.exists():
            found = sorted(arg.rglob("artifacts.json"))
            if not found:
                sys.exit(f"No artifacts.json under {arg}")
            cand = found[0]
        return cand
    if not arg.exists():
        sys.exit(f"No such file: {arg}")
    return arg


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(
        description="Turn a perturb-seq artifacts.json into a Google "
                    "Slides-ready .pptx deck.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Import: upload the .pptx to Drive and open with Google Slides, "
               "or use File > Import slides to pull individual slides into an "
               "existing deck (choose 'Keep original theme').",
    )
    p.add_argument("artifacts", type=Path,
                   help="path to artifacts.json, or the analysis directory")
    p.add_argument("-o", "--out", type=Path,
                   help="output .pptx (default: qc_slides.pptx beside artifacts.json)")
    p.add_argument("--title", type=str, help="deck title (default: run directory name)")
    p.add_argument("--subtitle", type=str, default="", help="deck subtitle")
    p.add_argument("--aspect", choices=("16:9", "4:3"), default="16:9")
    p.add_argument("--per-slide", type=int, choices=(1, 2, 3), default=2,
                   help="max half-width figures per slide (default 2)")
    p.add_argument("--pack", choices=("width", "aspect", "family"),
                   default="family",
                   help="'family' (default) groups a run of '<x> by <condition>' "
                        "figures of the same size into one labelled grid slide, "
                        "and otherwise behaves like 'aspect'; 'aspect' pairs up "
                        "anything that is not genuinely wide; 'width' mirrors "
                        "the HTML's half/full setting exactly")
    p.add_argument("--family-max", type=int, default=6, metavar="N",
                   help="most figures to put in one family grid (default 6; "
                        "1 disables family grouping)")
    p.add_argument("--pack-aspect", type=float, default=1.7,
                   help="with --pack aspect, figures wider than this ratio keep "
                        "a slide to themselves (default 1.7)")
    p.add_argument("--from-report", type=Path, metavar="QC_REPORT_HTML",
                   help="also read the built HTML report, to pick up the run "
                        "provenance table and the report title (the one thing "
                        "the pipeline never writes to artifacts.json)")
    p.add_argument("--tables", choices=("charts", "full", "none"),
                   default="charts",
                   help="'charts' (default) draws a table as a bar chart where "
                        "it has a label column and numbers, and drops it "
                        "otherwise; 'full' lays tables out as tables; 'none' "
                        "leaves them out entirely")
    p.add_argument("--no-tables", action="store_true",
                   help="alias for --tables none")
    p.add_argument("--no-tables-in", type=str, default="perturbation",
                   metavar="SECTIONS",
                   help="comma-separated sections to drop every table from "
                        "(default 'perturbation'; pass '' to keep them)")
    p.add_argument("--max-charts", type=int, default=3, metavar="N",
                   help="most table-derived charts to add per section "
                        "(default 3; 0 for no limit)")
    p.add_argument("--max-notes", type=int, default=4, metavar="N",
                   help="most method-note blocks to keep per section, ranked "
                        "poor > warn > info then by whether they quantify "
                        "something (default 4; 0 for no limit)")
    p.add_argument("--all-notes", action="store_true",
                   help="keep method notes that are pipeline bookkeeping rather "
                        "than findings (they are dropped by default; warn/poor "
                        "notes are always kept)")
    p.add_argument("--keep-all-tables", action="store_true",
                   help="keep tables whose numbers are already plotted (they "
                        "are dropped by default; see REDUNDANT_TABLES)")
    p.add_argument("--drop-tables", type=str, metavar="REGEX,REGEX",
                   help="also drop tables whose title matches any of these, "
                        "e.g. --drop-tables '^Knockdown by '")
    p.add_argument("--keep-tables", type=str, metavar="REGEX,REGEX",
                   help="protect tables whose title matches any of these, "
                        "overriding every drop rule")
    p.add_argument("--max-table-rows", type=int, default=11,
                   help="rows per table slide (default 12)")
    p.add_argument("--max-table-cols", type=int, default=5,
                   help="columns per table slide; wider tables are split into "
                        "column groups that repeat the first column (default 5, "
                        "which is what fits legibly at 12pt)")
    p.add_argument("--max-image-px", type=int, default=1800,
                   help="downscale figures above this long edge (0 disables; "
                        "needs Pillow)")
    p.add_argument("--banner-max", type=float, default=1.45,
                   help="section prose shorter than this (inches) shares the "
                        "first figure slide instead of getting its own "
                        "(default 1.45; 0 always gives it its own slide)")
    p.add_argument("--sections", type=str,
                   help="comma-separated section keys to include")
    p.add_argument("--skip-sections", type=str,
                   help="comma-separated section keys to drop")
    p.add_argument("--keep-duplicates", action="store_true",
                   help="keep figures that are byte-identical to one already in "
                        "the deck (they are dropped by default)")
    p.add_argument("--no-conclusions", action="store_true",
                   help="omit the blank 'Conclusions & next steps' slide that "
                        "closes the deck")
    p.add_argument("--conclusions", type=str, metavar="A|B|C",
                   help="pipe-separated headings for the closing slide "
                        "(default 'Conclusions|Open questions|Next steps')")
    p.add_argument("--appendix", action="store_true",
                   help="include the appendix section (configuration, "
                        "provenance, input bookkeeping); omitted by default")
    p.add_argument("--no-contents", action="store_true",
                   help="omit the contents slide")
    p.add_argument("-q", "--quiet", action="store_true")
    args = p.parse_args(argv)

    src = resolve_input(args.artifacts)
    out = args.out or (src.parent / "qc_slides.pptx")
    html_title, provenance = (None, [])
    if args.from_report:
        if not args.from_report.exists():
            sys.exit(f"No such report: {args.from_report}")
        html_title, provenance = parse_report_html(args.from_report)
    build_deck(
        src, out,
        title=args.title or html_title,
        subtitle=args.subtitle,
        aspect=args.aspect,
        per_slide=args.per_slide,
        table_mode="none" if args.no_tables else args.tables,
        max_table_rows=args.max_table_rows,
        max_table_cols=args.max_table_cols,
        drop_redundant_tables=not args.keep_all_tables,
        no_table_sections=[x.strip() for x in args.no_tables_in.split(",") if x.strip()],
        essential_notes_only=not args.all_notes,
        max_notes_per_section=args.max_notes,
        max_charts_per_section=args.max_charts,
        drop_tables=[x.strip() for x in args.drop_tables.split(",")] if args.drop_tables else (),
        keep_tables=[x.strip() for x in args.keep_tables.split(",")] if args.keep_tables else (),
        conclusions=() if args.no_conclusions else (
            [x.strip() for x in args.conclusions.split("|") if x.strip()]
            if args.conclusions
            else ("Conclusions", "Open questions", "Next steps")
        ),
        max_image_px=args.max_image_px,
        banner_max=args.banner_max,
        pack=args.pack,
        pack_aspect=args.pack_aspect,
        family_max=args.family_max,
        drop_duplicates=not args.keep_duplicates,
        provenance=provenance,
        only=[s.strip() for s in args.sections.split(",")] if args.sections else None,
        appendix=args.appendix,
        skip=[s.strip() for s in args.skip_sections.split(",")] if args.skip_sections else None,
        contents=not args.no_contents,
        quiet=args.quiet,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
